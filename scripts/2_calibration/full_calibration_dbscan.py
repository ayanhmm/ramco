#!/usr/bin/env python3
"""
full_calibration_dbscan.py

Enhanced with progress bars for long-running grid search.
"""

import os
import re
import cv2
import numpy as np
import pandas as pd
import itertools
from pptx import Presentation
from sklearn.cluster import DBSCAN
from sklearn.linear_model import LinearRegression
from tqdm import tqdm

# ─── CONFIGURATION ─────────────────────────────────────────────────────────────
# Paths to PPTX and image directories
WRAP_PPTX     = "data/ramco_images/Sheet stack with stretch wrap.pptx"
UNWRAP_PPTX   = "data/ramco_images/Sheet stack without stretch wrap.pptx"
WRAP_DIR      = "data/raw/wrap_images"
UNWRAP_DIR    = "data/raw/nowrap_images"
# Parameter ranges
PHT_THRESH_RANGE   = [50, 100, 150]
MIN_LINE_LEN_RANGE = [30, 50, 70]
MAX_LINE_GAP_RANGE = [3, 5, 8]
DBSCAN_EPS_RANGE   = [5, 8, 12]
IMAGE_EXTS         = ("jpg","jpeg","png")
# ────────────────────────────────────────────────────────────────────────────────

def extract_mapping_and_gt(pptx_path, image_dir):
    prs    = Presentation(pptx_path)
    slides = list(prs.slides)[1:]
    out = []
    for i, slide in enumerate(slides, start=2):
        idx = i - 1
        ts = getattr(slide.shapes, "title", None)
        if not ts or not getattr(ts, "has_text_frame", False):
            ts = next((s for s in slide.shapes if getattr(s, "has_text_frame", False)), None)
        caption = ts.text.strip() if ts and ts.has_text_frame else ""
        m = re.search(r"(\d+)\s*No", caption)
        if not m:
            continue
        gt = int(m.group(1))
        for ext in IMAGE_EXTS:
            name = f"image{idx}.{ext}"
            if os.path.exists(os.path.join(image_dir, name)):
                out.append({"image": name, "gt_count": gt})
                break
    return out

def preprocess_bw(path):
    img  = cv2.imread(path)
    gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
    gray = cv2.convertScaleAbs(gray, alpha=1.25, beta=0)
    sob  = cv2.Sobel(gray, cv2.CV_64F, dx=0, dy=1, ksize=3)
    abs_sob = np.uint8(np.absolute(sob) / (np.abs(sob).max()+1e-6) * 255)
    return cv2.Canny(abs_sob, 50, 150)

def raw_count_dbscan(bw, pht_thresh, min_len, max_gap, eps):
    segs = cv2.HoughLinesP(bw, 1, np.pi/180, pht_thresh,
                          minLineLength=min_len,
                          maxLineGap=max_gap)
    if segs is None:
        return 0
    mids = []
    for x1,y1,x2,y2 in segs[:,0]:
        ang = abs(np.degrees(np.arctan2(y2-y1, x2-x1)))
        if ang < 1 or abs(ang-180) < 1:
            mids.append([(y1+y2)/2.0])
    if not mids:
        return 0
    labels = DBSCAN(eps=eps, min_samples=1).fit(np.array(mids)).labels_
    return len([l for l in set(labels) if l != -1])

def main():
    wrap_map   = extract_mapping_and_gt(WRAP_PPTX, WRAP_DIR)
    unwrap_map = extract_mapping_and_gt(UNWRAP_PPTX, UNWRAP_DIR)

    combos = list(itertools.product(
        PHT_THRESH_RANGE,
        MIN_LINE_LEN_RANGE,
        MAX_LINE_GAP_RANGE,
        DBSCAN_EPS_RANGE
    ))

    best = {"mse": float("inf")}
    best_params = None

    print(f"Grid searching over {len(combos)} parameter combinations...")
    for pht, mn, mg, eps in tqdm(combos, desc="Grid Search", unit="combo"):
        errs = []
        for mapping, img_dir in [(wrap_map, WRAP_DIR), (unwrap_map, UNWRAP_DIR)]:
            for m in mapping:
                bw = preprocess_bw(os.path.join(img_dir, m["image"]))
                pred = raw_count_dbscan(bw, pht, mn, mg, eps)
                errs.append((pred - m["gt_count"])**2)
        mse = sum(errs)/len(errs)
        if mse < best["mse"]:
            best["mse"] = mse
            best_params = (pht, mn, mg, eps)

    pht, mn, mg, eps = best_params
    print(f"\nBest parameters: PHT={pht}, minLineLen={mn}, maxLineGap={mg}, eps={eps} (MSE={best['mse']:.2f})")

    records = []
    for mapping, img_dir in [(wrap_map, WRAP_DIR), (unwrap_map, UNWRAP_DIR)]:
        for m in mapping:
            bw  = preprocess_bw(os.path.join(img_dir, m["image"]))
            raw = raw_count_dbscan(bw, pht, mn, mg, eps)
            records.append({
                "image":    m["image"],
                "wrapped":  img_dir == WRAP_DIR,
                "gt_count": m["gt_count"],
                "pred_raw": raw
            })

    df = pd.DataFrame(records)
    reg = LinearRegression().fit(df[["pred_raw"]], df["gt_count"])
    df["pred_cal"] = reg.predict(df[["pred_raw"]]).round().astype(int)
    print(f"Post-regression MSE: {((df.pred_cal - df.gt_count)**2).mean():.2f}")

    df.to_csv("calibration_results_dbscan.csv", index=False)
    print("Saved calibration_results_dbscan.csv")

if __name__ == "__main__":
    main()

'''
python scripts/2_calibration/full_calibration_dbscan.py 
'''