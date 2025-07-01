#!/usr/bin/env python3
"""
full_calibration.py

1. Read both PPTX decks, skip slide 1, extract captions + map to image files.
2. Parse out ground-truth counts from captions ("NN No").
3. Run Sobel+Otsu+Hough ridge counter on each image.
4. Build a DataFrame of gt_count vs. pred_raw.
5. Grid-search Hough parameters to minimize MSE.
6. (Optional) Fit a linear regression for final bias correction.
7. Print summary and save calibration_results.csv.
"""

import os, re, cv2, numpy as np, pandas as pd
from pptx import Presentation
from sklearn.linear_model import LinearRegression

# ──────────────────────────────────────────────────────────────────────────────
# Paths to PPTX and image directories
WRAP_PPTX     = "data/ramco_images/Sheet stack with stretch wrap.pptx"
UNWRAP_PPTX   = "data/ramco_images/Sheet stack without stretch wrap.pptx"
WRAP_DIR      = "data/raw/wrap_images"
UNWRAP_DIR    = "data/raw/nowrap_images"
# ──────────────────────────────────────────────────────────────────────────────

# Initial Hough parameters (we’ll grid-search around these)
HOUGH_THRESHOLDS = [100, 150, 200]
MIN_GAPS         = [5, 10, 15]
ANGLE_TOL        = np.pi/180 * 5
IMAGE_EXTS       = ("jpg","jpeg","png")

def extract_mapping_and_gt(pptx_path, image_dir):
    """
    Skips slide 1. For slides 2…N:
      - slide_index → imageX where X = slide_index - 1
      - caption = first placeholder title or first text frame
      - gt_count = number before 'No' in caption
    Returns list of dicts: {image_name, gt_count}
    """
    prs    = Presentation(pptx_path)
    slides = list(prs.slides)[1:]  # skip the cover
    mappings = []
    for i, slide in enumerate(slides, start=2):
        img_idx = i - 1
        # 1. caption
        title_shape = getattr(slide.shapes, "title", None)
        if not title_shape or not getattr(title_shape, "has_text_frame", False):
            title_shape = next((s for s in slide.shapes if getattr(s, "has_text_frame", False)), None)
        caption = title_shape.text.strip() if title_shape and title_shape.has_text_frame else ""
        # 2. parse gt_count
        m = re.search(r"(\d+)\s*No", caption)
        if not m:
            continue
        gt = int(m.group(1))
        # 3. find corresponding image file
        image_name = None
        for ext in IMAGE_EXTS:
            candidate = f"image{img_idx}.{ext}"
            if os.path.exists(os.path.join(image_dir, candidate)):
                image_name = candidate
                break
        if image_name:
            mappings.append({"image": image_name, "gt_count": gt})
    return mappings

def preprocess_bw(path):
    img  = cv2.imread(path)
    gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
    sob  = cv2.Sobel(gray, cv2.CV_64F, dx=0, dy=1, ksize=3)
    abs_sob = np.uint8(cv2.normalize(np.abs(sob), None, 0,255,cv2.NORM_MINMAX))
    _, bw   = cv2.threshold(abs_sob, 0, 255, cv2.THRESH_BINARY+cv2.THRESH_OTSU)
    return bw

def count_hough(bw, thresh, min_gap):
    lines = cv2.HoughLines(bw,1,np.pi/180,thresh)
    if lines is None: return 0
    y0s = [rho/(np.sin(theta)+1e-6)
           for rho,theta in lines[:,0]
           if abs(theta)<ANGLE_TOL or abs(theta-np.pi)<ANGLE_TOL]
    if not y0s: return 0
    y0s.sort()
    clusters=[y0s[0]]
    for y in y0s[1:]:
        if abs(y-clusters[-1])>min_gap:
            clusters.append(y)
    return len(clusters)

def main():
    # 1. Gather wrapped & unwrapped mappings
    wrap_map   = extract_mapping_and_gt(WRAP_PPTX, WRAP_DIR)
    unwrap_map = extract_mapping_and_gt(UNWRAP_PPTX, UNWRAP_DIR)

    # 2. Build raw predictions DataFrame
    records = []
    for wrapped, img_dir, mapping in [
        (True, WRAP_DIR, wrap_map),
        (False, UNWRAP_DIR, unwrap_map)
    ]:
        for m in mapping:
            path = os.path.join(img_dir, m["image"])
            bw   = preprocess_bw(path)
            pred = count_hough(bw, HOUGH_THRESHOLDS[1], MIN_GAPS[1])
            records.append({
                "image":     m["image"],
                "wrapped":   wrapped,
                "gt_count":  m["gt_count"],
                "pred_raw":  pred
            })

    df = pd.DataFrame(records)
    print("\nInitial calibration data:")
    print(df, "\n")

    # 3. Grid search for best Hough params
    best = {"mse": float("inf")}
    for ht in HOUGH_THRESHOLDS:
        for mg in MIN_GAPS:
            errs = []
            for _, row in df.iterrows():
                bw   = preprocess_bw(os.path.join(
                    row.wrapped and WRAP_DIR or UNWRAP_DIR,
                    row.image
                ))
                p = count_hough(bw, ht, mg)
                errs.append((p - row.gt_count)**2)
            mse = sum(errs)/len(errs)
            if mse < best["mse"]:
                best = {"hough_thresh": ht, "min_gap": mg, "mse": mse}
    print("Best Hough params:", best, "\n")

    # 4. Optional linear correction
    X = df[["pred_raw"]].values
    y = df["gt_count"].values
    reg = LinearRegression().fit(X, y)
    df["pred_cal"] = reg.predict(X).round().astype(int)
    print("Post-regression MSE:", ((df.pred_cal - df.gt_count)**2).mean(), "\n")

    # 5. Save full calibration table
    df.to_csv("calibration_results.csv", index=False)
    print("Calibration results saved to calibration_results.csv")

if __name__=="__main__":
    main()

'''
python scripts/2_calibration/full_calibration.py 
'''