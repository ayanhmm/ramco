#!/usr/bin/env python3
"""
train_feature_regressors.py

– Reads slide‐titles → ground‐truth counts
– Extracts features per image: raw_count, total_length, total_edges
– Trains two RandomForestRegressor models (wrapped vs. unwrapped)
– Reports CV MSE and saves rf_wrapped.pkl & rf_unwrapped.pkl
"""

import os
import re
import cv2
import numpy as np
import pandas as pd
import pickle
from pptx import Presentation
from sklearn.ensemble import RandomForestRegressor
from sklearn.model_selection import cross_val_score, KFold
from sklearn.pipeline import make_pipeline
from sklearn.preprocessing import StandardScaler
from sklearn.cluster import DBSCAN

# ─── UPDATE THESE PATHS ─────────────────────────────────────────────────────────
WRAP_PPTX = "ramco_images/Sheet stack with stretch wrap.pptx"
NOWRAP_PPTX = "ramco_images/Sheet stack without stretch wrap.pptx"
WRAP_DIR = "ramco_images/wrap_images"
NOWRAP_DIR = "ramco_images/nowrap_images"
# ────────────────────────────────────────────────────────────────────────────────

IMAGE_EXTS = ("jpg","jpeg","png")

def extract_mapping(pptx_path, image_dir):
    prs    = Presentation(pptx_path)
    slides = list(prs.slides)[1:]  # skip cover slide
    records = []
    for i, slide in enumerate(slides, start=2):
        idx = i - 1
        ts = getattr(slide.shapes, "title", None)
        if not ts or not getattr(ts, "has_text_frame", False):
            ts = next((s for s in slide.shapes if getattr(s,"has_text_frame", False)), None)
        caption = ts.text.strip() if ts else ""
        m = re.search(r"(\d+)\s*No", caption)
        if not m:
            continue
        gt = int(m.group(1))
        for ext in IMAGE_EXTS:
            fname = f"image{idx}.{ext}"
            fpath = os.path.join(image_dir, fname)
            if os.path.exists(fpath):
                records.append({"path": fpath, "gt": gt})
                break
    return pd.DataFrame(records)

def extract_features(df):
    rows = []
    for _, r in df.iterrows():
        img = cv2.imread(r["path"])
        gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
        gray = cv2.convertScaleAbs(gray, alpha=1.25, beta=0)
        sob = cv2.Sobel(gray, cv2.CV_64F, 0, 1, 3)
        abs_sob = np.abs(sob)
        norm = np.uint8(abs_sob / (abs_sob.max()+1e-6) * 255)
        edges = cv2.Canny(norm, 50, 150)
        total_edges = edges.sum() / 255

        # Probabilistic Hough + DBSCAN
        segs = cv2.HoughLinesP(edges, 1, np.pi/180, 150, minLineLength=70, maxLineGap=8)
        raw_count = 0
        total_length = 0
        mids = []
        if segs is not None:
            for x1, y1, x2, y2 in segs[:,0]:
                angle = abs(np.degrees(np.arctan2(y2 - y1, x2 - x1)))
                if angle < 1 or abs(angle - 180) < 1:
                    total_length += np.hypot(x2 - x1, y2 - y1)
                    mids.append([(y1 + y2) / 2.0])
            if mids:
                labels = DBSCAN(eps=5, min_samples=1).fit(np.array(mids)).labels_
                raw_count = len([lbl for lbl in set(labels) if lbl != -1])

        rows.append({
            "raw_count": raw_count,
            "total_length": total_length,
            "total_edges": total_edges,
            "gt": r["gt"]
        })
    return pd.DataFrame(rows)

def train_and_save(df, name):
    X = df[["raw_count", "total_length", "total_edges"]].values
    y = df["gt"].values
    pipeline = make_pipeline(
        StandardScaler(),
        RandomForestRegressor(n_estimators=100, random_state=0)
    )
    cv = KFold(5, shuffle=True, random_state=0)
    mse = -cross_val_score(pipeline, X, y, cv=cv,
                           scoring="neg_mean_squared_error").mean()
    print(f"{name} CV MSE: {mse:.1f}")
    pipeline.fit(X, y)
    with open(f"rf_{name}.pkl", "wb") as f:
        pickle.dump(pipeline, f)
    df["pred"] = pipeline.predict(X)
    print(f"\n{name} predictions vs. GT:\n", df[["gt","pred"]])
    print(f"Saved model to rf_{name}.pkl\n")
    return pipeline

if __name__ == "__main__":
    df_wrap = extract_mapping(WRAP_PPTX, WRAP_DIR)
    df_nowrap = extract_mapping(NOWRAP_PPTX, NOWRAP_DIR)

    feat_wrap = extract_features(df_wrap)
    feat_nowrap = extract_features(df_nowrap)

    train_and_save(feat_wrap, "wrapped")
    train_and_save(feat_nowrap, "unwrapped")

