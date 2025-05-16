#!/usr/bin/env python3
"""
train_multimodal.py

– Auto-caption each image with BLIP
– Encode caption with CLIP text encoder
– Extract CV features: raw_count, total_length, total_edges
– Concatenate [CV_feats, CLIP_text_emb] → train regressor
– Shows progress bars and loads models just once.
"""

import os, re, cv2, numpy as np, pandas as pd, pickle, torch
from pptx import Presentation
from sklearn.ensemble import RandomForestRegressor
from sklearn.model_selection import cross_val_score, KFold
from sklearn.pipeline import make_pipeline
from sklearn.preprocessing import StandardScaler
from sklearn.cluster import DBSCAN
from tqdm import tqdm

# transformers imports
from transformers import BlipProcessor, BlipForConditionalGeneration
from transformers import CLIPProcessor, CLIPModel

# ─── CONFIG ─────────────────────────────────────────────────────────────────────
WRAP_PPTX    = "ramco_images/Sheet stack with stretch wrap.pptx"
NOWRAP_PPTX  = "ramco_images/Sheet stack without stretch wrap.pptx"
WRAP_DIR = "ramco_images/wrap_images"
NOWRAP_DIR = "ramco_images/nowrap_images"
IMAGE_EXTS   = ("jpg","jpeg","png")
# ────────────────────────────────────────────────────────────────────────────────

# 1) Load BLIP & CLIP once
print("Loading BLIP captioning model…")
blip_processor = BlipProcessor.from_pretrained("Salesforce/blip-image-captioning-base")
blip_model     = BlipForConditionalGeneration.from_pretrained("Salesforce/blip-image-captioning-base")
print("Loading CLIP text encoder…")
clip_processor = CLIPProcessor.from_pretrained("openai/clip-vit-base-patch32")
clip_model     = CLIPModel.from_pretrained("openai/clip-vit-base-patch32")
print("Models loaded.\n")

def auto_caption(image_path):
    img = cv2.imread(image_path)[..., ::-1]  # BGR → RGB
    inputs = blip_processor(images=img, return_tensors="pt")
    out    = blip_model.generate(**inputs, max_new_tokens=30)
    return blip_processor.decode(out[0], skip_special_tokens=True)

def clip_text_embed(caption):
    inputs = clip_processor(text=[caption], return_tensors="pt", padding=True)
    with torch.no_grad():
        txt_feats = clip_model.get_text_features(**inputs)
    return txt_feats[0].cpu().numpy()  # (512,)

def extract_mapping(pptx_path, image_dir):
    prs    = Presentation(pptx_path)
    slides = list(prs.slides)[1:]
    recs   = []
    for i, slide in enumerate(slides, start=2):
        idx = i-1
        ts = getattr(slide.shapes, "title", None)
        if not ts or not getattr(ts, "has_text_frame", False):
            ts = next((s for s in slide.shapes if getattr(s,"has_text_frame",False)), None)
        caption = ts.text.strip() if ts else ""
        m = re.search(r"(\d+)\s*No", caption)
        if not m:
            continue
        gt = int(m.group(1))
        for ext in IMAGE_EXTS:
            name = f"image{idx}.{ext}"
            path = os.path.join(image_dir, name)
            if os.path.exists(path):
                recs.append({"path":path, "gt":gt})
                break
    return pd.DataFrame(recs)

def extract_features_and_text(df, label):
    rows = []
    print(f"Extracting features & captions for {label} ({len(df)} images)…")
    for _, r in tqdm(df.iterrows(), total=len(df), desc=label):
        # CV features
        img = cv2.imread(r["path"])
        gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
        gray = cv2.convertScaleAbs(gray, alpha=1.25, beta=0)
        sob = cv2.Sobel(gray, cv2.CV_64F, 0,1,3)
        norm = np.uint8(np.abs(sob) / (np.abs(sob).max()+1e-6) * 255)
        edges = cv2.Canny(norm, 50, 150)
        total_edges = edges.sum()/255

        segs = cv2.HoughLinesP(edges,1,np.pi/180,150,minLineLength=70,maxLineGap=8)
        raw = length = 0
        mids = []
        if segs is not None:
            for x1,y1,x2,y2 in segs[:,0]:
                ang = abs(np.degrees(np.arctan2(y2-y1, x2-x1)))
                if ang<1 or abs(ang-180)<1:
                    length += np.hypot(x2-x1,y2-y1)
                    mids.append([(y1+y2)/2.])
            if mids:
                labels = DBSCAN(eps=5, min_samples=1).fit(np.array(mids)).labels_
                raw = len([l for l in set(labels) if l!=-1])

        # Text features
        cap = auto_caption(r["path"])
        txt_emb = clip_text_embed(cap)  # (512,)

        # Combine
        feat = np.concatenate([
            np.array([raw, length, total_edges], dtype=np.float32),
            txt_emb
        ])
        rows.append({"features":feat, "gt":r["gt"]})
    return rows

def train_multimodal(rows, name):
    X = np.stack([r["features"] for r in rows])
    y = np.array([r["gt"] for r in rows])
    pipe = make_pipeline(StandardScaler(),
                         RandomForestRegressor(n_estimators=100, random_state=0))
    cv = KFold(5, shuffle=True, random_state=0)
    mse = -cross_val_score(pipe, X, y, cv=cv,
                           scoring="neg_mean_squared_error").mean()
    print(f"{name} CV MSE: {mse:.1f}")
    pipe.fit(X,y)
    pickle.dump(pipe, open(f"rf_{name}_multi.pkl","wb"))
    print(f"Saved rf_{name}_multi.pkl\n")
    return pipe

if __name__ == "__main__":
    # Wrapped
    df_wrap = extract_mapping(WRAP_PPTX, WRAP_DIR)
    rows_w = extract_features_and_text(df_wrap, "Wrapped")
    train_multimodal(rows_w, "wrapped")

    # Unwrapped
    df_nowrap = extract_mapping(NOWRAP_PPTX, NOWRAP_DIR)
    rows_n = extract_features_and_text(df_nowrap, "Unwrapped")
    train_multimodal(rows_n, "unwrapped")

