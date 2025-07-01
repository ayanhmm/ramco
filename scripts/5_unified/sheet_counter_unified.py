#!/usr/bin/env python3
"""
sheet_counter_unified.py

Pure-model hybrid:
- Always extract [raw, length, edges, layer] CV features.
- Always extract CLIP text embedding from slide title text (if any).
- PCA-reduce text features and run the RF regressor to predict counts.
"""

import os, glob, pickle, cv2, numpy as np, torch
from pptx import Presentation
from tqdm import tqdm
from sklearn.cluster import DBSCAN
from transformers import CLIPProcessor, CLIPModel

# ─── CONFIG ───────────────────────────────────────────────────────────────────
RF_WRAP_BUNDLE   = "./models/rf_wrapped_multi_pca_layers.pkl"
RF_NOWRAP_BUNDLE = "./models/rf_unwrapped_multi_pca_layers.pkl"
CLIP_MODEL_NAME  = "openai/clip-vit-base-patch32"
# ───────────────────────────────────────────────────────────────────────────────

# Load trained bundles
wrap_bundle   = pickle.load(open(RF_WRAP_BUNDLE,   "rb"))
nowrap_bundle = pickle.load(open(RF_NOWRAP_BUNDLE, "rb"))
pca_wrap,   model_wrap   = wrap_bundle["pca"],   wrap_bundle["model"]
pca_nowrap, model_nowrap = nowrap_bundle["pca"], nowrap_bundle["model"]

# Initialize CLIP text encoder
clip_proc  = CLIPProcessor.from_pretrained(CLIP_MODEL_NAME)
clip_model = CLIPModel.from_pretrained(CLIP_MODEL_NAME)
clip_model.eval()

def load_title_texts(pptx_path, img_dir):
    """Map image filename → slide title text."""
    texts = {}
    if pptx_path and os.path.exists(pptx_path):
        # print(f"Loading titles from {pptx_path}...")
        prs    = Presentation(pptx_path)
        slides = list(prs.slides)[1:]
        # print(f"Found {len(slides)} slides in {pptx_path}.")
        
        for i, slide in enumerate(slides, start=2):
            # print(f"Processing slide {i}...")
            ts = getattr(slide.shapes, "title", None)
            # print(f"  Title shape: {ts}")
            if not ts or not ts.has_text_frame:
                ts = next((s for s in slide.shapes if getattr(s,"has_text_frame",False)), None)
            txt = ts.text.strip() if ts else ""
            idx = i - 1
            # print(f"  Slide title: '{i}' '{txt}'")
            
            # fn = f"image{idx}.jpg"
            # texts[fn] = txt
            
            
            for ext in ("jpg","jpeg","png"):
                fn = f"image{idx}.{ext}"
                fn2= os.path.join(img_dir, fn)
                # print(f"  Checking for image: {fn2}")
                if os.path.exists(fn2):
                    texts[fn] = txt
                    break
    # print(texts)
    return texts

def extract_features(path, title_map):
    """Compute CV feats + CLIP-text-PCA feats; return combined feature vector."""
    print(path)
    img = cv2.imread(path)
    gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
    gray = cv2.convertScaleAbs(gray, alpha=1.25, beta=0)

    # ── CV horizontal edges → raw_count & length
    sob_y = cv2.Sobel(gray, cv2.CV_64F, 0,1,3)
    norm_y = np.uint8(np.abs(sob_y)/(np.abs(sob_y).max()+1e-6)*255)
    ed_y = cv2.Canny(norm_y, 50,150)
    segs_h = cv2.HoughLinesP(ed_y,1,np.pi/180,150,70,8)
    raw = length = 0
    mids = []
    if segs_h is not None:
        for x1,y1,x2,y2 in segs_h[:,0]:
            ang = abs(np.degrees(np.arctan2(y2-y1, x2-x1)))
            if ang<1 or abs(ang-180)<1:
                length += np.hypot(x2-x1, y2-y1)
                mids.append([(y1+y2)/2.])
        if mids:
            labels_h = DBSCAN(eps=5, min_samples=1).fit(np.array(mids)).labels_
            raw = len([l for l in set(labels_h) if l!=-1])
    total_edges = ed_y.sum()/255

    # ── CV vertical edges → layer_count
    sob_x = cv2.Sobel(gray, cv2.CV_64F, 1,0,3)
    norm_x = np.uint8(np.abs(sob_x)/(np.abs(sob_x).max()+1e-6)*255)
    ed_x = cv2.Canny(norm_x, 50,150)
    segs_v = cv2.HoughLinesP(ed_x,1,np.pi/180,100,30,5)
    layer_count = 0
    midsx = []
    if segs_v is not None:
        for x1,y1,x2,y2 in segs_v[:,0]:
            ang = abs(np.degrees(np.arctan2(y2-y1, x2-x1)))
            if abs(ang-90)<5:
                midsx.append([(x1+x2)/2.])
        if midsx:
            labels_v = DBSCAN(eps=10, min_samples=1).fit(np.array(midsx)).labels_
            layer_count = len([l for l in set(labels_v) if l!=-1])

    cv_feats = [raw, length, total_edges, layer_count]

    # ── CLIP text → embed → PCA
    fn = os.path.basename(path)
    # print(fn)
    txt = title_map.get(fn, "")
    print(f"Extracting features for {fn} with title '{txt}'")
    inputs = clip_proc(text=[txt], return_tensors="pt", padding=True)
    # print(inputs)
    with torch.no_grad():
        text_emb = clip_model.get_text_features(**inputs)[0].cpu().numpy()

    # choose correct PCA by folder name
    folder = os.path.basename(os.path.dirname(path))
    if folder == "wrap_images":
        text_pca = pca_wrap.transform([text_emb])[0]
    else:
        text_pca = pca_nowrap.transform([text_emb])[0]

    return np.hstack([cv_feats, text_pca])

def process(folder, title_map, bundle, label):
    pca = bundle["pca"]
    model = bundle["model"]

    imgs = sorted(glob.glob(os.path.join(folder,"*.jpg")) +
                glob.glob(os.path.join(folder,"*.jpeg")) +
                glob.glob(os.path.join(folder,"*.png")))
    # print(imgs)
    subtotal = 0
    print(f"\n--- {label} ({len(imgs)} images) ---")
    
    
    for p in tqdm(imgs, desc=label):
        feat = extract_features(p, title_map)
        # print(feat)
        cnt  = int(round(model.predict([feat])[0]))
        subtotal += cnt
        print(f"{os.path.basename(p):25s} → {cnt:3d}")
    print(f"Subtotal {label}: {subtotal}")
    return subtotal

if __name__ == "__main__":
    import argparse
    parser = argparse.ArgumentParser(__doc__)
    parser.add_argument("--wrap_pptx",   required=True, help="PPTX for wrapped titles")
    parser.add_argument("--nowrap_pptx", required=True, help="PPTX for unwrapped titles")
    parser.add_argument("--wrap_dir",    required=True,  help="Folder of wrapped images")
    parser.add_argument("--nowrap_dir",  required=True,  help="Folder of unwrapped images")
    args = parser.parse_args()

    nowrap_titles = []
    wrap_titles = []
    # wrap_titles   = load_title_texts(args.wrap_pptx,   args.wrap_dir)
    nowrap_titles = load_title_texts(args.nowrap_pptx, args.nowrap_dir)
    print(nowrap_titles)
    
    print(f"Loaded {len(wrap_titles)} wrapped titles and {len(nowrap_titles)} unwrapped titles.")

    tot_u = 0  
    tot_w = 0
    # tot_w = process(args.wrap_dir,   wrap_titles,   wrap_bundle,   "Wrapped")
    tot_u = process(args.nowrap_dir, nowrap_titles, nowrap_bundle, "Unwrapped")
    

    print("\n" + "="*40)
    print(f"GRAND TOTAL: {tot_w + tot_u}")
    print("="*40)

'''
python scripts/5_unified/sheet_counter_unified.py \
    --wrap_pptx  data/ramco_images/"Sheet stack with stretch wrap.pptx" \
    --nowrap_pptx data/ramco_images/"Sheet stack without stretch wrap.pptx" \
    --wrap_dir   data/raw/wrap_images \
    --nowrap_dir data/raw/nowrap_images
'''