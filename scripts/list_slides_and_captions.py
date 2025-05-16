#!/usr/bin/env python3
"""
list_slides_and_captions.py

Reads both PPTX files, extracts each slide’s title and
matches it to the corresponding image filename (imageX.ext).
Skips the first (cover) slide, so slide 2 → image1, slide 3 → image2, etc.
Prints two tables: Wrapped and Unwrapped.
"""

import os
from pptx import Presentation
import pandas as pd

# Paths
WRAP_PPTX    = "ramco_images/Sheet stack with stretch wrap.pptx"
UNWRAP_PPTX  = "ramco_images/Sheet stack without stretch wrap.pptx"
WRAP_DIR     = "ramco_images/wrap_images"
UNWRAP_DIR   = "ramco_images/nowrap_images"  # corrected folder name

IMAGE_EXTS = ("jpg", "jpeg", "png")

def extract_info(pptx_path, image_dir):
    """
    Returns a list of dicts with:
      slide_number (actual PPT slide index),
      image_index (slide_number - 1),
      image_name (or '(missing)'),
      title         (or '(no title)').
    Skips the very first slide.
    """
    prs = Presentation(pptx_path)
    slides = list(prs.slides)
    info = []

    # Start at slide 2 (index 1 in zero-based list), map to image1
    for slide_number, slide in enumerate(slides[1:], start=2):
        image_index = slide_number - 1

        # 1) Try built-in title placeholder
        title_shape = getattr(slide.shapes, "title", None)
        if not title_shape or not getattr(title_shape, "has_text_frame", False):
            # 2) Fallback: first shape with text
            title_shape = next(
                (sh for sh in slide.shapes if getattr(sh, "has_text_frame", False)),
                None
            )

        title = title_shape.text.strip() if title_shape and title_shape.has_text_frame else "(no title)"

        # Find the matching image file (jpg, jpeg, or png)
        image_name = "(missing)"
        for ext in IMAGE_EXTS:
            candidate = f"image{image_index}.{ext}"
            if os.path.exists(os.path.join(image_dir, candidate)):
                image_name = candidate
                break

        info.append({
            "slide_number": slide_number,
            "image_index":  image_index,
            "image_name":   image_name,
            "title":        title
        })

    return info


if __name__ == "__main__":
    # Wrapped
    wrapped = extract_info(WRAP_PPTX, WRAP_DIR)
    df_wrap = pd.DataFrame(wrapped)
    print("\n=== Wrapped Slides (slide2 → image1, etc.) ===")
    print(df_wrap.to_string(index=False))

    # Unwrapped
    unwrapped = extract_info(UNWRAP_PPTX, UNWRAP_DIR)
    df_unwrap = pd.DataFrame(unwrapped)
    print("\n=== Unwrapped Slides (slide2 → image1, etc.) ===")
    print(df_unwrap.to_string(index=False))

