#!/usr/bin/env python3
"""
train_multimodal_pca.py

– Same as before but caps PCA components to #samples
"""

import os, re, json, cv2, numpy as np, pandas as pd, pickle, torch
from pptx import Presentation
from tqdm import tqdm
from sklearn.cluster import DBSCAN
from sklearn.decomposition import PCA
from sklearn.ensemble import RandomForestRegressor
from sklearn.model_selection import cross_val_score, KFold
from sklearn.pipeline import make_pipeline
from sklearn.preprocessing import StandardScaler
from transformers import CLIPProcessor, CLIPModel

# ─── CONFIG ───────────────────────────────────────────────────────────────────
CAP_JSON         = "captions.json"
WRAP_PPTX        = "ramco_images/Sheet stack with stretch wrap.pptx"
NOWRAP_PPTX      = "ramco_images/Sheet stack without stretch wrap.pptx"
WRAP_DIR         = "ramco_images/wrap_images"
NOWRAP_DIR       = "ramco_images/nowrap_images"
IMAGE_EXTS       = ("jpg","jpeg","png")
DESIRED_PCA_COMP = 30
# ────────────────────────────────────────────────────────────────────────────────

# Load cached captions & CLIP once
captions  = json.load(open(CAP_JSON))
clip_proc = CLIPProcessor.from_pretrained("openai/clip-vit-base-patch32")
clip_model= CLIPModel.from_pretrained("openai/clip-vit-base-patch32")
clip_model.eval()

def get_text_emb(fn):
    cap = captions.get(fn, "")
    inp = clip_proc(text=[cap], return_tensors="pt", padding=True)
    with torch.no_grad():
        return clip_model.get_text_features(**inp)[0].cpu().numpy()

def extract_mapping(pptx, img_dir):
    prs    = Presentation(pptx)
    slides = list(prs.slides)[1:]
    recs   = []
    for i, slide in enumerate(slides, start=2):
        idx = i-1
        ts  = getattr(slide.shapes, "title", None)
        if not ts or not ts.has_text_frame:
            ts = next((s for s in slide.shapes if s.has_text_frame), None)
        text = ts.text.strip() if ts else ""
        m = re.search(r"(\d+)\s*No", text)
        if not m: continue
        gt = int(m.group(1))
        for ext in IMAGE_EXTS:
            fn = f"image{idx}.{ext}"
            p  = os.path.join(img_dir, fn)
            if os.path.exists(p):
                recs.append({"fn":fn, "path":p, "gt":gt})
                break
    return pd.DataFrame(recs)

def extract_all_features(df, label):
    cv_feats, txt_embs, gts = [], [], []
    for _, r in tqdm(df.iterrows(), total=len(df), desc=label):
        img = cv2.imread(r["path"])
        gray= cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
        gray= cv2.convertScaleAbs(gray, alpha=1.25, beta=0)
        sob = cv2.Sobel(gray, cv2.CV_64F, 0,1,3)
        norm= np.uint8(np.abs(sob)/(np.abs(sob).max()+1e-6)*255)
        edges= cv2.Canny(norm, 50,150)
        total_edges = edges.sum()/255

        segs = cv2.HoughLinesP(edges,1,np.pi/180,150,minLineLength=70,maxLineGap=8)
        raw=length=0; mids=[]
        if segs is not None:
            for x1,y1,x2,y2 in segs[:,0]:
                ang = abs(np.degrees(np.arctan2(y2-y1, x2-x1)))
                if ang<1 or abs(ang-180)<1:
                    length += np.hypot(x2-x1, y2-y1)
                    mids.append([(y1+y2)/2.])
            if mids:
                lbl = DBSCAN(eps=5, min_samples=1).fit(np.array(mids)).labels_
                raw = len([l for l in set(lbl) if l!=-1])

        cv_feats.append([raw, length, total_edges])
        txt_embs.append(get_text_emb(r["fn"]))
        gts.append(r["gt"])

    return np.array(cv_feats), np.stack(txt_embs), np.array(gts)

def train_with_pca(cv_feats, txt_embs, gts, name):
    n_samples, n_feats = txt_embs.shape
    n_pca = min(DESIRED_PCA_COMP, n_samples)
    pca = PCA(n_components=n_pca, random_state=0)
    txt_pca = pca.fit_transform(txt_embs)
    print(f"{name}: PCA comps={n_pca}, var_ratio_sum={pca.explained_variance_ratio_.sum():.3f}")

    X = np.hstack([cv_feats, txt_pca])
    y = gts

    pipe = make_pipeline(StandardScaler(),
                         RandomForestRegressor(n_estimators=100, random_state=0))
    cv   = KFold(5, shuffle=True, random_state=0)
    mse  = -cross_val_score(pipe, X, y, cv=cv,
                            scoring="neg_mean_squared_error").mean()
    print(f"{name} CV MSE: {mse:.1f}")

    pipe.fit(X, y)
    pickle.dump({"pca":pca, "model":pipe},
                open(f"rf_{name}_multi_pca.pkl","wb"))
    print(f"Saved rf_{name}_multi_pca.pkl\n")

if __name__=="__main__":
    dfw = extract_mapping(WRAP_PPTX, WRAP_DIR)
    cvw, txtw, gw = extract_all_features(dfw, "Wrapped")
    train_with_pca(cvw, txtw, gw, "wrapped")

    dfn = extract_mapping(NOWRAP_PPTX, NOWRAP_DIR)
    cvn, txtn, gn = extract_all_features(dfn, "Unwrapped")
    train_with_pca(cvn, txtn, gn, "unwrapped")

