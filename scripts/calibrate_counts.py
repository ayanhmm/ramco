#!/usr/bin/env python3
"""
calibrate_counts.py

1. Extract ground-truth counts from two PPTX files.
2. Compare against raw predictions from sheet_counter_hough pipeline.
3. Parameter sweep for best Hough + clustering params.
4. (Optional) Fit linear regression to correct residual bias.
"""

import os
import re
import cv2
import numpy as np
import pandas as pd
from pptx import Presentation
from sklearn.linear_model import LinearRegression

# Paths to PPTX and image directories
WRAP_PPTX     = "ramco_images/Sheet stack with stretch wrap.pptx"
UNWRAP_PPTX   = "ramco_images/Sheet stack without stretch wrap.pptx"
WRAP_DIR      = "ramco_images/wrap_images"
UNWRAP_DIR    = "ramco_images/unwrap_images"

# Hough parameters default
DEFAULT_HOUGH_THRESH = 150
DEFAULT_MIN_GAP      = 10
ANGLE_TOLERANCE      = np.pi/180 * 5
IMAGE_EXTS           = ("jpg", "jpeg", "png")

def parse_ground_truth(pptx_path, image_dir):
    """
    Extract slide titles containing counts 'NN No' → dict(image_name → count).
    Chooses the correct extension by checking files in image_dir.
    """
    prs = Presentation(pptx_path)
    gt = {}
    for idx, slide in enumerate(prs.slides, start=1):
        shape = next((s for s in slide.shapes if s.has_text_frame), None)
        if not shape:
            continue
        text = shape.text.strip()
        m = re.search(r"(\d+)\s*No", text)
        if not m:
            continue
        count = int(m.group(1))
        # find the actual image file with any supported extension
        for ext in IMAGE_EXTS:
            name = f"image{idx}.{ext}"
            if os.path.exists(os.path.join(image_dir, name)):
                gt[name] = count
                break
    return gt

def preprocess_bw(path):
    """Read image, Sobel for horizontal edges, Otsu threshold to binary."""
    img  = cv2.imread(path)
    gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
    sob  = cv2.Sobel(gray, cv2.CV_64F, dx=0, dy=1, ksize=3)
    abs_sob = np.uint8(cv2.normalize(np.abs(sob), None, 0, 255, cv2.NORM_MINMAX))
    _, bw   = cv2.threshold(abs_sob, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)
    return bw

def count_sheets(bw, hough_thresh, min_gap):
    """Detect near-horizontal lines via Hough and cluster by y-intercept."""
    lines = cv2.HoughLines(bw, 1, np.pi/180, hough_thresh)
    if lines is None:
        return 0
    y0s = []
    for rho, theta in lines[:,0]:
        if abs(theta) < ANGLE_TOLERANCE or abs(theta - np.pi) < ANGLE_TOLERANCE:
            y0 = rho / (np.sin(theta) + 1e-6)
            y0s.append(y0)
    if not y0s:
        return 0
    y0s = sorted(y0s)
    clusters = [y0s[0]]
    for y in y0s[1:]:
        if abs(y - clusters[-1]) > min_gap:
            clusters.append(y)
    return len(clusters)

if __name__ == "__main__":
    # 1. Extract ground truth for wrapped & unwrapped
    wrapped_gt   = parse_ground_truth(WRAP_PPTX, WRAP_DIR)
    unwrapped_gt = parse_ground_truth(UNWRAP_PPTX, UNWRAP_DIR)

    # 2. Build calibration DataFrame
    records = []
    for name, gt_count in wrapped_gt.items():
        path = os.path.join(WRAP_DIR, name)
        bw   = preprocess_bw(path)
        pred = count_sheets(bw, DEFAULT_HOUGH_THRESH, DEFAULT_MIN_GAP)
        records.append({
            "image":     name,
            "wrapped":   True,
            "gt_count":  gt_count,
            "pred_raw":  pred
        })
    for name, gt_count in unwrapped_gt.items():
        path = os.path.join(UNWRAP_DIR, name)
        bw   = preprocess_bw(path)
        pred = count_sheets(bw, DEFAULT_HOUGH_THRESH, DEFAULT_MIN_GAP)
        records.append({
            "image":     name,
            "wrapped":   False,
            "gt_count":  gt_count,
            "pred_raw":  pred
        })

    df = pd.DataFrame(records)
    print("Calibration data:\n", df)

    # 3. Parameter sweep: find best (hough_thresh, min_gap)
    best = {"err": float("inf")}
    for ht in [100, 150, 200]:
        for mg in [5, 10, 15]:
            errs = []
            for _, row in df.iterrows():
                dir_ = WRAP_DIR if row.wrapped else UNWRAP_DIR
                bw   = preprocess_bw(os.path.join(dir_, row.image))
                p    = count_sheets(bw, ht, mg)
                errs.append((p - row.gt_count) ** 2)
            mse = sum(errs) / len(errs)
            if mse < best["err"]:
                best = {"hough_thresh": ht, "min_gap": mg, "err": mse}

    print("\nBest parameters:", best)

    # 4. (Optional) Fit linear regression to remove residual bias
    X = df[["pred_raw"]].values
    y = df["gt_count"].values
    reg = LinearRegression().fit(X, y)
    df["pred_calibrated"] = reg.predict(X).round().astype(int)
    mse_cal = ((df["pred_calibrated"] - df["gt_count"])**2).mean()
    print("Post-regression MSE:", mse_cal)

    # 5. Save results
    df.to_csv("calibration_results.csv", index=False)

