#!/usr/bin/env python3
"""
retrain_cv_only_models.py

Re-extract CV-only features and retrains two RF pipelines:
  • rf_wrapped_multi.pkl
  • rf_unwrapped_multi.pkl (if any data)
"""

import os, re, pickle, cv2, numpy as np, pandas as pd
from pptx import Presentation
from sklearn.ensemble import RandomForestRegressor
from sklearn.pipeline import Pipeline
from sklearn.preprocessing import StandardScaler
from sklearn.cluster import DBSCAN

# ─── CONFIG ─────────────────────────────────────────────────────────────────────
WRAP_PPTX      = "ramco_images/Sheet stack with stretch wrap.pptx"
NOWRAP_PPTX    = "ramco_images/Sheet stack without stretch wrap.pptx"
WRAP_DIR       = "ramco_images/wrap_images"
NOWRAP_DIR     = "ramco_images/nowrap_images"
OUT_WRAP_PKL   = "rf_wrapped_multi.pkl"
OUT_NOWRAP_PKL = "rf_unwrapped_multi.pkl"
RE_BUNDLES     = re.compile(r"(\d+)[×xX*]\s*(\d+)")
# ────────────────────────────────────────────────────────────────────────────────

# up near the top, add:
RE_BUNDLES = re.compile(r"(\d+)[×xX*]\s*(\d+)")
RE_NO      = re.compile(r"(\d+)\s*No", re.IGNORECASE)

def parse_gt(pptx_path, img_dir):
    prs    = Presentation(pptx_path)
    slides = list(prs.slides)[1:]  # skip title slide
    mapping = {}
    for idx, slide in enumerate(slides, start=2):
        # grab title text
        title_shape = getattr(slide.shapes, "title", None)
        if not title_shape or not title_shape.has_text_frame:
            title_shape = next((s for s in slide.shapes if getattr(s,"has_text_frame",False)), None)
        txt = title_shape.text.strip() if title_shape else ""

        # 1) try the bundles×per form
        m = RE_BUNDLES.search(txt)
        if m:
            bundles, per = map(int, m.groups())
            total = bundles * per
        else:
            # 2) fallback to “NN No’s”
            m2 = RE_NO.search(txt)
            if not m2:
                continue
            total = int(m2.group(1))

        # map slide→image filename
        img_idx = idx - 1   # slide2 → image1
        for ext in ("jpg","jpeg","png"):
            fn = f"image{img_idx}.{ext}"
            if os.path.exists(os.path.join(img_dir, fn)):
                mapping[fn] = total
                break

    return mapping

def compute_cv_feats(path):
    img = cv2.imread(path)
    gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)

    # horizontal edges → raw_count & length
    sob_y = cv2.Sobel(gray, cv2.CV_64F, 0, 1, 3)
    norm_y = np.uint8(np.abs(sob_y)/(np.abs(sob_y).max()+1e-6)*255)
    edges_y = cv2.Canny(norm_y, 50, 150)
    segs_h = cv2.HoughLinesP(edges_y, 1, np.pi/180, 100, 30, 10)
    mids_y, total_length = [], 0
    if segs_h is not None:
        for x1,y1,x2,y2 in segs_h[:,0]:
            ang = abs(np.degrees(np.arctan2(y2-y1, x2-x1)))
            if ang < 5 or abs(ang-180) < 5:
                mids_y.append([(y1+y2)/2.])
                total_length += np.hypot(x2-x1, y2-y1)
    raw_count = len(set(DBSCAN(eps=5, min_samples=1)
                        .fit(np.array(mids_y) if mids_y else [[0]])
                        .labels_))

    # total edge pixels
    total_edges = float(edges_y.sum()) / 255.0

    # vertical edges → layer count
    sob_x = cv2.Sobel(gray, cv2.CV_64F, 1, 0, 3)
    norm_x = np.uint8(np.abs(sob_x)/(np.abs(sob_x).max()+1e-6)*255)
    edges_x = cv2.Canny(norm_x, 50, 150)
    segs_v = cv2.HoughLinesP(edges_x, 1, np.pi/180, 80, 20, 5)
    mids_x = []
    if segs_v is not None:
        for x1,y1,x2,y2 in segs_v[:,0]:
            ang = abs(np.degrees(np.arctan2(y2-y1, x2-x1)))
            if abs(ang-90) < 5:
                mids_x.append([(x1+x2)/2.])
    layer_count = len(set(DBSCAN(eps=5, min_samples=1)
                          .fit(np.array(mids_x) if mids_x else [[0]])
                          .labels_))

    return [raw_count, total_length, total_edges, layer_count]

def build_df(gt_map, folder):
    records = []
    for fn, gt in gt_map.items():
        feats = compute_cv_feats(os.path.join(folder, fn))
        records.append({
            "image": fn, "gt": gt,
            "raw": feats[0],
            "length": feats[1],
            "edges": feats[2],
            "layers": feats[3]
        })
    return pd.DataFrame(records)

def train_and_save(df, out_pkl, label):
    X = df[["raw","length","edges","layers"]].values
    y = df["gt"].values
    pipe = Pipeline([
        ("scale", StandardScaler()),
        ("rf", RandomForestRegressor(n_estimators=100, random_state=42))
    ])
    pipe.fit(X, y)
    mse = ((pipe.predict(X) - y)**2).mean()
    print(f"{label}: Trained on {len(df)} samples; MSE = {mse:.1f}")
    with open(out_pkl, "wb") as f:
        pickle.dump(pipe, f)
    print(f"→ saved {out_pkl}\n")

if __name__ == "__main__":
    print("Parsing GT counts…")
    wrap_map   = parse_gt(WRAP_PPTX,  WRAP_DIR)
    nowrap_map = parse_gt(NOWRAP_PPTX, NOWRAP_DIR)

    dfw = build_df(wrap_map, WRAP_DIR)
    print(f"\nWrapped DF: {dfw.shape}")
    if not dfw.empty:
        train_and_save(dfw, OUT_WRAP_PKL, "Wrapped")
    else:
        print("No Wrapped samples—skipping.\n")

    dfu = build_df(nowrap_map, NOWRAP_DIR)
    print(f"Unwrapped DF: {dfu.shape}")
    if not dfu.empty:
        train_and_save(dfu, OUT_NOWRAP_PKL, "Unwrapped")
    else:
        print("No Unwrapped samples—skipping.\n")




