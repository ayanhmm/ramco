#!/usr/bin/env python3
"""
train_multimodal_cached.py

– Loads precomputed captions from captions.json
– Embeds each caption with CLIP
– Extracts CV features
– Trains multimodal regressor
Improved: import torch for no_grad
"""

import os
import re
import json
import cv2
import numpy as np
import pandas as pd
import pickle
import torch
from pptx import Presentation
from tqdm import tqdm
from sklearn.cluster import DBSCAN
from sklearn.ensemble import RandomForestRegressor
from sklearn.model_selection import cross_val_score, KFold
from sklearn.pipeline import make_pipeline
from sklearn.preprocessing import StandardScaler
from transformers import CLIPProcessor, CLIPModel

# ─── CONFIG ───────────────────────────────────────────────────────────────────
CAP_JSON         = "data/previews/captions.json"
WRAP_PPTX     = "data/ramco_images/Sheet stack with stretch wrap.pptx"
NOWRAP_PPTX   = "data/ramco_images/Sheet stack without stretch wrap.pptx"
WRAP_DIR      = "data/raw/wrap_images"
NOWRAP_DIR    = "data/raw/nowrap_images"
IMAGE_EXTS = ("jpg","jpeg","png")
# ────────────────────────────────────────────────────────────────────────────────

# Load cached captions
with open(CAP_JSON, "r") as f:
    captions = json.load(f)

# Load CLIP once
clip_proc  = CLIPProcessor.from_pretrained("openai/clip-vit-base-patch32")
clip_model = CLIPModel.from_pretrained("openai/clip-vit-base-patch32")
clip_model.eval()

def get_text_emb(fn):
    caption = captions.get(fn, "")
    inputs  = clip_proc(text=[caption], return_tensors="pt", padding=True)
    with torch.no_grad():
        feat = clip_model.get_text_features(**inputs)
    return feat[0].cpu().numpy()

def extract_mapping(pptx, img_dir):
    prs    = Presentation(pptx)
    slides = list(prs.slides)[1:]
    recs   = []
    for i, slide in enumerate(slides, start=2):
        idx = i-1
        ts  = getattr(slide.shapes, "title", None)
        if not ts or not getattr(ts, "has_text_frame", False):
            ts = next((s for s in slide.shapes if getattr(s,"has_text_frame",False)), None)
        title = ts.text.strip() if ts else ""
        m = re.search(r"(\d+)\s*No", title)
        if not m:
            continue
        gt = int(m.group(1))
        for ext in IMAGE_EXTS:
            fn = f"image{idx}.{ext}"
            path = os.path.join(img_dir, fn)
            if os.path.exists(path):
                recs.append({"fn":fn, "path":path, "gt":gt})
                break
    return pd.DataFrame(recs)

def extract_features(df, label):
    rows = []
    for _, r in tqdm(df.iterrows(), total=len(df), desc=label):
        img = cv2.imread(r["path"])
        gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
        gray = cv2.convertScaleAbs(gray, alpha=1.25, beta=0)
        sob = cv2.Sobel(gray, cv2.CV_64F, 0,1,3)
        norm = np.uint8(np.abs(sob)/(np.abs(sob).max()+1e-6)*255)
        edges = cv2.Canny(norm, 50,150)
        total_edges = edges.sum()/255

        segs = cv2.HoughLinesP(edges,1,np.pi/180,150,minLineLength=70,maxLineGap=8)
        raw = length = 0
        mids = []
        if segs is not None:
            for x1,y1,x2,y2 in segs[:,0]:
                ang = abs(np.degrees(np.arctan2(y2-y1,x2-x1)))
                if ang<1 or abs(ang-180)<1:
                    length += np.hypot(x2-x1,y2-y1)
                    mids.append([(y1+y2)/2.])
            if mids:
                labels = DBSCAN(eps=5,min_samples=1).fit(np.array(mids)).labels_
                raw = len([l for l in set(labels) if l!=-1])

        txt_emb = get_text_emb(r["fn"])
        feat = np.concatenate([[raw, length, total_edges], txt_emb])
        rows.append({"feat":feat, "gt":r["gt"]})
    return rows

def train(rows, name):
    X = np.stack([r["feat"] for r in rows])
    y = np.array([r["gt"] for r in rows])
    pipe = make_pipeline(StandardScaler(), RandomForestRegressor(n_estimators=100, random_state=0))
    cv   = KFold(5, shuffle=True, random_state=0)
    mse  = -cross_val_score(pipe, X, y, cv=cv, scoring="neg_mean_squared_error").mean()
    print(f"{name} CV MSE: {mse:.1f}")
    pipe.fit(X, y)
    pickle.dump(pipe, open(f"rf_{name}_multi.pkl", "wb"))
    print(f"Saved rf_{name}_multi.pkl\n")

if __name__=="__main__":
    dfw   = extract_mapping(WRAP_PPTX, WRAP_DIR)
    rowsw = extract_features(dfw, "Wrapped")
    train(rowsw, "wrapped")

    dfn   = extract_mapping(NOWRAP_PPTX, NOWRAP_DIR)
    rowsn = extract_features(dfn, "Unwrapped")
    train(rowsn, "unwrapped")

'''
python scripts/4_multimodal/train_multimodal_cached.py 
'''