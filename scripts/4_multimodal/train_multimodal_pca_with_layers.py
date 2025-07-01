#!/usr/bin/env python3
"""
train_multimodal_pca_with_layers.py

– Loads cached captions (CLIP emb)
– Extracts CV features + layer_count per image
– Applies PCA to compress text embeddings
– Concatenates [raw_count, total_length, total_edges, layer_count, txt_pca]
– Trains RandomForest and saves rf_{wrapped,unwrapped}_pca_layers.pkl
"""

import os, re, json, cv2, numpy as np, pandas as pd, pickle, torch
from pptx import Presentation
from tqdm import tqdm
from sklearn.cluster import DBSCAN
from sklearn.decomposition import PCA
from sklearn.ensemble import RandomForestRegressor
from sklearn.model_selection import cross_val_score, KFold
from sklearn.pipeline import make_pipeline
from sklearn.preprocessing import StandardScaler
from transformers import CLIPProcessor, CLIPModel

# ─── CONFIG ───────────────────────────────────────────────────────────────────
CAP_JSON         = "data/previews/captions.json"
WRAP_PPTX     = "data/ramco_images/Sheet stack with stretch wrap.pptx"
NOWRAP_PPTX   = "data/ramco_images/Sheet stack without stretch wrap.pptx"
WRAP_DIR      = "data/raw/wrap_images"
NOWRAP_DIR    = "data/raw/nowrap_images"
IMAGE_EXTS       = ("jpg","jpeg","png")
DESIRED_PCA_COMP = 30
# ────────────────────────────────────────────────────────────────────────────────

captions  = json.load(open(CAP_JSON))
clip_proc = CLIPProcessor.from_pretrained("openai/clip-vit-base-patch32")
clip_model= CLIPModel.from_pretrained("openai/clip-vit-base-patch32")
clip_model.eval()

def get_text_emb(fn):
    inp = clip_proc(text=[captions.get(fn,"")], return_tensors="pt", padding=True)
    with torch.no_grad():
        return clip_model.get_text_features(**inp)[0].cpu().numpy()

def extract_mapping(pptx, img_dir):
    prs    = Presentation(pptx)
    slides = list(prs.slides)[1:]
    recs   = []
    for i, slide in enumerate(slides, start=2):
        idx = i-1
        ts = getattr(slide.shapes, "title", None)
        if not ts or not ts.has_text_frame:
            ts = next((s for s in slide.shapes if s.has_text_frame), None)
        title = ts.text.strip() if ts else ""
        # print(f"Processing slide {i} with title: '{title}'")
        m = re.search(r"(\d+)\s*No", title)
        if not m: continue
        gt = int(m.group(1))
        for ext in IMAGE_EXTS:
            fn = f"image{idx}.{ext}"
            p  = os.path.join(img_dir, fn)
            if os.path.exists(p):
                recs.append({"fn":fn, "path":p, "gt":gt})
                break
    return pd.DataFrame(recs)

def extract_all_features(df, label, vis_dir="data/previews/4_Train_multimodal_PCA_with_layers"):
    os.makedirs(vis_dir, exist_ok=True) if vis_dir else None
    cv_feats, txt_embs, gts = [], [], []

    for _, r in tqdm(df.iterrows(), total=len(df), desc=label):
        img = cv2.imread(r["path"])
        display_img = img.copy()
        gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
        gray = cv2.convertScaleAbs(gray, alpha=1.25, beta=0)

        # # --- Horizontal edge detection ---
        # sob_y = cv2.Sobel(gray, cv2.CV_64F, dx=0, dy=1, ksize=3)
        # norm_y = np.uint8(np.abs(sob_y) / (np.abs(sob_y).max() + 1e-6) * 255)
        # edges_y = cv2.Canny(norm_y, 50, 150)
        # segs_h = cv2.HoughLinesP(edges_y, 1, np.pi/180, 150, minLineLength=70, maxLineGap=8)
        # --- Enhanced Horizontal edge detection ---
        sob_y = cv2.Sobel(gray, cv2.CV_64F, dx=0, dy=1, ksize=5)  # increase ksize for stronger gradients
        abs_sob_y = np.abs(sob_y)
        norm_y = np.uint8(abs_sob_y / (abs_sob_y.max() + 1e-6) * 255)

        # More sensitive Canny thresholds
        edges_y = cv2.Canny(norm_y, 20, 80)  # Lower thresholds catch more edges

        # More permissive Hough Transform
        segs_h = cv2.HoughLinesP(edges_y, 1, np.pi/180, 100, minLineLength=40, maxLineGap=12)

        raw = length = 0
        mids_h = []
        if segs_h is not None:
            for x1, y1, x2, y2 in segs_h[:, 0]:
                dx, dy = x2 - x1, y2 - y1
                ang = abs(np.degrees(np.arctan2(dy, dx)))

                # Accept lines within ±5° of horizontal (more aggressive than <1°)
                # And ensure they are reasonably long and not vertical-ish
                if ang < 5 or abs(ang - 180) < 5:
                    line_len = np.hypot(dx, dy)

                    # Only accept lines longer than a tighter threshold
                    if line_len > 100 and abs(dy) < 5:
                        length += line_len
                        mids_h.append([(y1 + y2) / 2.])
                        cv2.line(display_img, (x1, y1), (x2, y2), (0, 255, 255), 1)  # yellow

            if mids_h:
                labels_h = DBSCAN(eps=5, min_samples=1).fit(np.array(mids_h)).labels_
                raw = len([l for l in set(labels_h) if l != -1])


        total_edges = edges_y.sum() / 255

        # --- Vertical edge detection ---
        sob_x = cv2.Sobel(gray, cv2.CV_64F, dx=1, dy=0, ksize=3)
        norm_x = np.uint8(np.abs(sob_x) / (np.abs(sob_x).max() + 1e-6) * 255)
        edges_x = cv2.Canny(norm_x, 50, 150)
        segs_v = cv2.HoughLinesP(edges_x, 1, np.pi/180, 100, minLineLength=30, maxLineGap=5)

        layer_count = 0
        mids_v = []
        if segs_v is not None:
            for x1, y1, x2, y2 in segs_v[:, 0]:
                ang = abs(np.degrees(np.arctan2(y2 - y1, x2 - x1)))
                if abs(ang - 90) < 5:
                    mids_v.append([(x1 + x2) / 2.])
                    cv2.line(display_img, (x1, y1), (x2, y2), (255, 255, 0), 1)  # cyan

            if mids_v:
                labels_v = DBSCAN(eps=10, min_samples=1).fit(np.array(mids_v)).labels_
                layer_count = len([l for l in set(labels_v) if l != -1])

        # Add overlay text
        cv2.putText(display_img, f"raw: {raw}, len: {int(length)}, edges: {int(total_edges)}, layers: {layer_count}",
                    (10, 25), cv2.FONT_HERSHEY_SIMPLEX, 0.6, (255, 0, 0), 2)

        # Save visualization
        if vis_dir:
            out_path = os.path.join(vis_dir, os.path.basename(r["path"]))
            cv2.imwrite(out_path, display_img)

        cv_feats.append([raw, length, total_edges, layer_count])
        txt_embs.append(get_text_emb(r["fn"]))
        gts.append(r["gt"])

    return np.array(cv_feats), np.stack(txt_embs), np.array(gts)


def train_with_pca(cv_feats, txt_embs, gts, name):
    # print(cv_feats)
    # print(gts)
    n_samples = txt_embs.shape[0]
    n_pca = min(DESIRED_PCA_COMP, n_samples)
    pca = PCA(n_components=n_pca, random_state=0)
    txt_pca = pca.fit_transform(txt_embs)
    print(f"{name}: PCA comps={n_pca}, var_ratio={pca.explained_variance_ratio_.sum():.3f}")

    X = np.hstack([cv_feats, txt_pca])
    y = gts

    pipe = make_pipeline(StandardScaler(), RandomForestRegressor(n_estimators=100, random_state=0))
    cv   = KFold(5, shuffle=True, random_state=0)
    mse  = -cross_val_score(pipe, X, y, cv=cv, scoring="neg_mean_squared_error").mean()
    print(f"{name} CV MSE: {mse:.1f}")

    pipe.fit(X, y)
    pickle.dump({"pca":pca, "model":pipe}, open(f"rf_{name}_multi_pca_layers.pkl","wb"))
    print(f"Saved rf_{name}_multi_pca_layers.pkl\n")

if __name__=="__main__":
    # dfw = extract_mapping(WRAP_PPTX, WRAP_DIR)
    # cvw, txtw, gw = extract_all_features(dfw, "Wrapped")
    # train_with_pca(cvw, txtw, gw, "wrapped")

    dfn = extract_mapping(NOWRAP_PPTX, NOWRAP_DIR)
    cvn, txtn, gn = extract_all_features(dfn, "Unwrapped")
    train_with_pca(cvn, txtn, gn, "unwrapped")

'''
python scripts/4_multimodal/train_multimodal_pca_with_layers.py 
'''